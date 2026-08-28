"""Arabic typography defaults for HerWork deliverables.

Cairo is the house font for anything that contains Arabic — the office
suite default (Calibri) falls back to an ugly Arabic rendering. These
helpers set Cairo on docx / pptx / pdf output, including the
complex-script font slot that Arabic text is actually shaped from.

For docx and pptx only the font NAME is embedded — the operating system
resolves it, so those work anywhere Cairo (or a fallback) is installed.
PDF generation via reportlab embeds the font file itself, so
``register_pdf_font`` locates the TTF across Linux/macOS/Windows font
directories and tells you to install Cairo if it can't.

Usage:
    from arabic_style import style_docx, style_pptx, register_pdf_font, FONT

    style_docx(doc)                      # after building the Document
    style_pptx(prs)                      # after building the Presentation
    register_pdf_font()                  # then canvas.setFont(FONT, size)
"""
import os
from pathlib import Path

FONT = "Cairo"

# Styles are enumerated from the document rather than listed here: a fixed
# tuple silently missed whatever the template happened to use (Heading 4-6,
# Caption, Quote, table styles), producing a report whose H4s render in the
# fallback font while its body is Cairo. These substrings mark the styles to
# LEAVE ALONE — monospace is load-bearing for code and Cairo would destroy
# the alignment that makes it readable.
_DOCX_KEEP_MONOSPACE = ("code", "macro", "preformatted", "plain text")


def _is_monospace_style(name) -> bool:
    if not name:
        return False
    lowered = str(name).lower()
    return any(hint in lowered for hint in _DOCX_KEEP_MONOSPACE)


def _font_dirs():
    """Font directories across Linux, macOS, and Windows."""
    home = Path.home()
    dirs = [
        home / ".local/share/fonts",          # Linux (user)
        home / ".fonts",                      # Linux (legacy user)
        Path("/usr/local/share/fonts"),       # Linux (local system)
        Path("/usr/share/fonts"),             # Linux (system)
        home / "Library/Fonts",               # macOS (user)
        Path("/Library/Fonts"),               # macOS (system)
    ]
    local_appdata = os.environ.get("LOCALAPPDATA")
    if local_appdata:                         # Windows 10+ per-user fonts
        dirs.append(Path(local_appdata) / "Microsoft/Windows/Fonts")
    windir = os.environ.get("WINDIR", r"C:\Windows")
    dirs.append(Path(windir) / "Fonts")       # Windows system fonts
    return dirs


def find_font_ttf(family: str = FONT) -> str:
    """Locate the TTF for an installed font family, cross-platform.

    Exact-family files (``Cairo-*.ttf``) win over lookalikes
    (``CairoPlay-*.ttf``). Raises FileNotFoundError with an install hint
    when the family isn't installed."""
    for base in _font_dirs():
        if not base.is_dir():
            continue
        hits = sorted(
            base.rglob(f"{family}*.ttf"),
            key=lambda p: (not p.name.startswith(f"{family}-"), len(p.name)),
        )
        if hits:
            return str(hits[0])
    raise FileNotFoundError(
        f"{family} font not found. Install it first, e.g. from "
        f"https://fonts.google.com/specimen/{family.replace(' ', '+')} "
        f"(or your OS package manager), then retry."
    )


def _docx_set_rfonts(element, qn) -> None:
    """Point an element's rPr at Cairo in the Latin AND complex-script slots."""
    rfonts = element.get_or_add_rPr().get_or_add_rFonts()
    for attr in ("w:ascii", "w:hAnsi", "w:cs"):
        rfonts.set(qn(attr), FONT)


def _docx_bodies(doc):
    """Yield every XML body that can hold text: document + headers/footers.

    Headers and footers live in separate parts, so a run in a running head
    is invisible to a walk of the document body alone. Only parts this
    document actually defines are touched — reading ``header._element`` on
    an inherited header would CREATE an empty one, so linked-to-previous
    headers are skipped rather than materialized.
    """
    yield doc.element.body
    for section in doc.sections:
        for attr in (
            "header", "footer",
            "first_page_header", "first_page_footer",
            "even_page_header", "even_page_footer",
        ):
            part = getattr(section, attr, None)          # older python-docx
            if part is None:
                continue
            try:
                if part.is_linked_to_previous:
                    continue
                element = part._element
            except (AttributeError, ValueError):
                continue
            if element is not None:
                yield element


def style_docx(doc) -> None:
    """Apply Cairo to every style AND every run in a python-docx Document.

    Sets the Latin (ascii/hAnsi) AND complex-script (cs) slots — Arabic is
    shaped from the cs slot, so setting only ``font.name`` leaves Arabic on
    the fallback font.

    Styling styles alone is not enough: a run carrying DIRECT formatting
    (its own ``w:rFonts``, which is what most docx-generating code and
    Word-authored templates emit) overrides its style, so a style-only pass
    returns success while the deliverable comes out half-Cairo,
    half-Calibri. Every run is therefore visited too — across the body,
    tables at any nesting depth, text boxes, and headers/footers.

    Runs in monospace-ish styles (code, macros, preformatted) keep their
    font: there the alignment is the point.
    """
    from docx.oxml.ns import qn

    # 1. Styles — covers text that carries no direct formatting, plus any
    #    style-driven text added to the document after this call.
    for style in doc.styles:
        name = getattr(style, "name", None)
        if _is_monospace_style(name):
            continue
        try:
            style.font.name = FONT
        except (AttributeError, NotImplementedError):
            pass  # table/numbering styles expose no .font — the XML below still applies
        try:
            _docx_set_rfonts(style.element, qn)
        except (AttributeError, TypeError):
            continue

    # 2. Direct run formatting — the half that actually decides the render.
    for body in _docx_bodies(doc):
        for paragraph in body.iter(qn("w:p")):
            pstyle = paragraph.find(qn("w:pPr") + "/" + qn("w:pStyle"))
            if pstyle is not None and _is_monospace_style(pstyle.get(qn("w:val"))):
                continue
            for run in paragraph.iter(qn("w:r")):
                _docx_set_rfonts(run, qn)


def _pptx_run_font(run) -> None:
    from pptx.oxml.ns import qn

    run.font.name = FONT  # sets a:latin
    rPr = run._r.get_or_add_rPr()
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = rPr.makeelement(qn("a:cs"), {})
        rPr.append(cs)
    cs.set("typeface", FONT)


def _pptx_frame_font(text_frame) -> None:
    for para in text_frame.paragraphs:
        for run in para.runs:
            _pptx_run_font(run)


def _pptx_walk_shapes(shapes) -> None:
    """Font every run under ``shapes``, descending into grouped shapes.

    A group shape has no ``text_frame`` of its own, so text inside it —
    common in real decks, and groups nest — is invisible to a flat pass
    over ``slide.shapes``. Recursion is depth-unbounded for that reason.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _pptx_walk_shapes(shape.shapes)
            continue
        if shape.has_text_frame:
            _pptx_frame_font(shape.text_frame)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    _pptx_frame_font(cell.text_frame)


def style_pptx(prs) -> None:
    """Set Cairo on every text run in a python-pptx Presentation.

    Covers table cells, grouped shapes at any nesting depth, and speaker
    notes — notes ship with the deliverable and render in the same fallback
    font as anything else when left alone.
    """
    for slide in prs.slides:
        _pptx_walk_shapes(slide.shapes)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            if notes is not None:
                _pptx_frame_font(notes)


def register_pdf_font() -> str:
    """Register Cairo with reportlab; returns the font name to setFont()."""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    pdfmetrics.registerFont(TTFont(FONT, find_font_ttf()))
    return FONT


def shape_arabic(text: str, font: str = FONT) -> str:
    """Reshape + reorder Arabic for a PDF canvas, with a glyph-coverage net.

    reportlab has no text shaper: it draws codepoint by codepoint, so Arabic
    must be pre-shaped into Presentation Forms-B (U+FE70-U+FEFF) and reordered
    to visual order before it is drawn. That is what ``arabic_reshaper`` and
    the bidi algorithm do.

    The trap is that a font is NOT required to carry that block. Cairo's
    variable build covers 89 of the 144 forms — everything a shaping engine
    would ever ask it for through GSUB, but not the standalone codepoints.
    Isolated alef (U+FE8D) and isolated teh (U+FE95) are among the missing,
    which is why "المبيعات" came out of a plain reshape+bidi pipeline as
    "▯لمبيعا▯": reportlab asked for a codepoint the font never claimed and
    drew .notdef, silently, in a finished deliverable.

    So every shaped character is checked against the registered font's own
    cmap and, when absent, folded back to its canonical letter (NFKC:
    U+FE8D -> U+0627). The base letters are all present, and for the forms
    that go missing the two glyphs are visually identical anyway — an
    isolated alef IS an alef.

    Call this instead of using ``arabic_reshaper``/``get_display`` directly,
    and only AFTER ``register_pdf_font()`` — the coverage check reads the
    registered face.
    """
    import unicodedata

    import arabic_reshaper
    from bidi.algorithm import get_display
    from reportlab.pdfbase import pdfmetrics

    visual = get_display(arabic_reshaper.reshape(text))

    try:
        supported = pdfmetrics.getFont(font).face.charToGlyph
    except Exception:
        # Font not registered, or a face that exposes no cmap (Type 1). The
        # shaped text is still the best answer available; don't fail the
        # document over a coverage check we can't run.
        return visual

    return "".join(
        character
        if ord(character) in supported
        else unicodedata.normalize("NFKC", character)
        for character in visual
    )
