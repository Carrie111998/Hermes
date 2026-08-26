"""Cairo must reach the text that actually renders, not just the styles.

A style-only pass reports success while a deliverable comes out half-Cairo
and half-Calibri: a run carrying DIRECT formatting (its own ``w:rFonts``,
which is what most docx-generating code and Word templates emit) overrides
its style. The pptx side had the mirror-image gap — grouped shapes have no
``text_frame``, so text inside them was skipped entirely.

These tests assert the rendered-font contract on real Document/Presentation
objects, since that is the only thing a reader of the file sees.
"""
import importlib.util
import sys
from pathlib import Path

import pytest

_SCRIPT = (
    Path(__file__).resolve().parents[2]
    / "optional-skills/productivity/herwork/scripts/arabic_style.py"
)


@pytest.fixture(scope="module")
def arabic_style():
    if not _SCRIPT.is_file():  # pragma: no cover - skill not checked out
        pytest.skip("herwork skill not present")
    spec = importlib.util.spec_from_file_location("herwork_arabic_style", _SCRIPT)
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def _docx_cs(run):
    from docx.oxml.ns import qn

    rfonts = run._r.find(qn("w:rPr") + "/" + qn("w:rFonts"))
    return None if rfonts is None else rfonts.get(qn("w:cs"))


def _pptx_cs(run):
    from pptx.oxml.ns import qn

    cs = run._r.find(qn("a:rPr") + "/" + qn("a:cs"))
    return None if cs is None else cs.get("typeface")


def test_style_docx_reaches_direct_run_formatting(arabic_style):
    """The half-Cairo/half-Calibri bug: styles alone never win over a run."""
    pytest.importorskip("docx")
    from docx import Document

    font = arabic_style.FONT
    doc = Document()

    direct = doc.add_paragraph().add_run("مرحبا")
    direct.font.name = "Calibri"  # direct formatting beats any style

    heading4 = doc.add_heading("عنوان رابع", level=4)  # beyond the old Heading 1-3
    caption = doc.add_paragraph("شرح الصورة", style="Caption")

    table = doc.add_table(rows=1, cols=1)
    cell_run = table.cell(0, 0).paragraphs[0].add_run("خلية")
    cell_run.font.name = "Calibri"
    nested = table.cell(0, 0).add_table(rows=1, cols=1)
    nested_run = nested.cell(0, 0).paragraphs[0].add_run("متداخل")
    nested_run.font.name = "Calibri"

    header_run = doc.sections[0].header.paragraphs[0].add_run("ترويسة")
    header_run.font.name = "Calibri"

    arabic_style.style_docx(doc)

    assert _docx_cs(direct) == font, "direct run formatting must be overridden"
    assert all(_docx_cs(r) == font for r in heading4.runs), "Heading 4 must be covered"
    assert all(_docx_cs(r) == font for r in caption.runs), "Caption must be covered"
    assert _docx_cs(cell_run) == font
    assert _docx_cs(nested_run) == font, "tables nest — the walk must too"
    assert _docx_cs(header_run) == font, "headers are a separate part"


def test_style_docx_leaves_monospace_styles_and_inherited_parts_alone(arabic_style):
    """Two things Cairo must not touch: code, and headers we never defined."""
    pytest.importorskip("docx")
    from docx import Document
    from docx.enum.style import WD_STYLE_TYPE

    doc = Document()
    doc.styles.add_style("Source Code", WD_STYLE_TYPE.PARAGRAPH)
    code = doc.add_paragraph("x = 1", style="Source Code").add_run("y = 2")
    code.font.name = "Courier New"
    normal = doc.add_paragraph().add_run("نص عادي")

    arabic_style.style_docx(doc)

    assert _docx_cs(code) != arabic_style.FONT, "monospace alignment is load-bearing"
    assert _docx_cs(normal) == arabic_style.FONT
    # Reading an inherited header/footer would CREATE an empty part; the walk
    # must skip linked-to-previous parts rather than materialize them.
    assert doc.sections[0].footer.is_linked_to_previous


def test_style_pptx_descends_into_groups_and_notes(arabic_style):
    """Grouped shapes hold no text frame of their own — and groups nest."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    font = arabic_style.FONT
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    plain = box.text_frame.paragraphs[0].add_run()
    plain.text = "نص"

    group = slide.shapes.add_group_shape()
    grouped_box = group.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
    grouped = grouped_box.text_frame.paragraphs[0].add_run()
    grouped.text = "داخل مجموعة"

    inner_group = group.shapes.add_group_shape()
    deep_box = inner_group.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
    deep = deep_box.text_frame.paragraphs[0].add_run()
    deep.text = "مجموعة داخل مجموعة"

    note = slide.notes_slide.notes_text_frame.paragraphs[0].add_run()
    note.text = "ملاحظات المتحدث"

    arabic_style.style_pptx(prs)

    assert _pptx_cs(plain) == font
    assert _pptx_cs(grouped) == font, "text inside a group must be reached"
    assert _pptx_cs(deep) == font, "recursion must be depth-unbounded"
    assert _pptx_cs(note) == font, "notes ship with the deliverable"
